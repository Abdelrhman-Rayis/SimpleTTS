"""Munfath (منفذ) — Mada Innovation Award pitch deck.

Simpler, startup-style layout: one idea per slide, generous whitespace,
big typography, real product screenshots embedded.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT     = Path("/Users/rayis/SimpleTTS")
ASSETS   = ROOT / "deck_assets"
OUT      = ROOT / "munfath_mada_pitch.pptx"

# Palette — pulled from the actual product so the deck matches the app.
CREAM    = RGBColor(0xF7, 0xEF, 0xE2)
PAPER    = RGBColor(0xFF, 0xFD, 0xF8)
INK      = RGBColor(0x10, 0x10, 0x10)
INK_2    = RGBColor(0x3A, 0x33, 0x2A)
MUTED    = RGBColor(0x6D, 0x62, 0x57)
SOFT     = RGBColor(0x96, 0x8C, 0x7F)
ACCENT   = RGBColor(0xE8, 0x5D, 0x2C)
ACCENT_2 = RGBColor(0xFF, 0xB1, 0x7A)
TEAL     = RGBColor(0x24, 0x5C, 0x69)

# Arabic font — falls back gracefully if not installed.
FONT_AR  = "Tajawal"
FONT_NUM = "SF Pro Display"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ── helpers ──────────────────────────────────────────────────────────────────

def bg(slide, color=CREAM):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s

def text(slide, x, y, w, h, content, *,
         size=18, bold=False, color=INK, font=FONT_AR,
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP, rtl=True,
         line_spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top  = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    if rtl:
        p._p.get_or_add_pPr().set('rtl', '1')
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = content
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb

def label(slide, x, y, text_str, *, color=ACCENT, size=10):
    """Small uppercase-style label on a slide (top tag)."""
    return text(slide, x, y, Inches(4), Inches(0.4), text_str,
                size=size, bold=True, color=color, align=PP_ALIGN.RIGHT)

def line(slide, x, y, w, h, color=INK, weight=1.0):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    return ln

def number_chip(slide, x, y, n, color=ACCENT):
    """Big numeral chip for step lists."""
    return text(slide, x, y, Inches(1.2), Inches(1.4), str(n),
                size=64, bold=True, color=color, font=FONT_NUM,
                align=PP_ALIGN.LEFT, rtl=False)

def img(slide, path, x, y, w=None, h=None, shadow=True):
    """Place an image with optional drop-shadow card."""
    p = ASSETS / path if not str(path).startswith("/") else Path(path)
    if not p.exists():
        return None
    with Image.open(p) as im:
        iw, ih = im.size
    if w and not h:
        h = int(w * ih / iw)
    elif h and not w:
        w = int(h * iw / ih)
    if shadow:
        # subtle shadow card behind
        sx = x + Inches(0.08); sy = y + Inches(0.12)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, sx, sy, w, h)
        card.adjustments[0] = 0.025
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x00,0x00,0x00)
        card.fill.transparency = 0
        card.line.fill.background()
        # We can't easily blur, so use a low-opacity dark via theme color hack.
        # Simpler: skip soft blur, just keep an outline-free dark behind image.
        # We'll instead just rely on the image itself looking framed.
        slide.shapes._spTree.remove(card._element)  # remove the rough shadow
    pic = slide.shapes.add_picture(str(p), x, y, width=w, height=h)
    return pic

def page_num(slide, n, total):
    text(slide, SW - Inches(1.3), SH - Inches(0.55), Inches(1), Inches(0.4),
         f"{n} / {total}", size=10, color=SOFT, align=PP_ALIGN.RIGHT, rtl=False)
    text(slide, Inches(0.5), SH - Inches(0.55), Inches(3.5), Inches(0.4),
         "منفذ  ·  mnfz.tech", size=10, color=SOFT, align=PP_ALIGN.LEFT, rtl=False)

TOTAL = 10

# ── Slide 1 — Title ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
# tiny brand stripe
line(s, Inches(0.7), Inches(0.7), Inches(0.3), Inches(0.06), color=ACCENT)
text(s, Inches(0.7), Inches(0.6), Inches(6), Inches(0.4),
     "MUNFATH · 2026", size=11, bold=True, color=ACCENT, align=PP_ALIGN.LEFT, rtl=False)
# big logo word
text(s, Inches(0), Inches(2.6), SW, Inches(2),
     "مُنفذ", size=180, bold=True, color=INK, align=PP_ALIGN.CENTER)
text(s, Inches(0), Inches(5.0), SW, Inches(0.6),
     "اقرأ بصوتٍ طبيعي. باللغة العربية. على جهازك.", size=22, color=MUTED,
     align=PP_ALIGN.CENTER)
text(s, Inches(0), SH - Inches(0.8), SW, Inches(0.4),
     "Mada Innovation Award 2026", size=11, color=SOFT,
     align=PP_ALIGN.CENTER, rtl=False)

# ── Slide 2 — The pain ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
label(s, SW - Inches(2.5), Inches(0.7), "المشكلة", color=ACCENT, size=11)
line(s, SW - Inches(2.8), Inches(1.1), Inches(0.3), Inches(0.04), color=ACCENT)
text(s, Inches(1), Inches(2.4), Inches(11.3), Inches(2.2),
     "ملايين القراء العرب\nبلا أداة قراءة ميسّرة.", size=64, bold=True, color=INK,
     align=PP_ALIGN.RIGHT, line_spacing=1.05)
text(s, Inches(1), Inches(5.0), Inches(11.3), Inches(1.5),
     "الأدوات العالمية مكلفة، وأصوات العربية فيها آلية وغير مفهومة، وكثير من ملفات الـ PDF غير قابلة للوصول أصلاً.",
     size=20, color=MUTED, align=PP_ALIGN.RIGHT, line_spacing=1.55)
page_num(s, 2, TOTAL)

# ── Slide 3 — The solution (big screenshot) ───────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, PAPER)
label(s, SW - Inches(2.5), Inches(0.55), "الحل", color=ACCENT, size=11)
text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0),
     "منفذ يفتح كل PDF بصوت طبيعي.", size=40, bold=True, color=INK,
     align=PP_ALIGN.RIGHT)
img(s, "home_light.png", Inches(1.2), Inches(2.1), w=Inches(11))
page_num(s, 3, TOTAL)

# ── Slide 4 — How it works ────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
label(s, SW - Inches(2.8), Inches(0.7), "كيف يعمل", color=ACCENT, size=11)
line(s, SW - Inches(3.1), Inches(1.1), Inches(0.3), Inches(0.04), color=ACCENT)
text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.8),
     "ثلاث خطوات.", size=44, bold=True, color=INK, align=PP_ALIGN.RIGHT)

steps = [
    ("١", "اختر ملف PDF",     "اسحب الملف أو اختره — يُفهرس منفذ النص تلقائياً مع الحفاظ على ترتيب القراءة العربي."),
    ("٢", "اضغط لتسمع",       "اضغط على أي فقرة لتُنطق فوراً بصوت عربي عصبي طبيعي، مع تمييز كل كلمة أثناء النطق."),
    ("٣", "تابع من حيث توقفت", "يحفظ منفذ موضع قراءتك تلقائياً، وزر متابعة فوري في مكتبتك الشخصية."),
]
col_w = Inches(3.85)
gap   = Inches(0.25)
total_w = col_w * 3 + gap * 2
start_x = (SW - total_w) // 2
for i, (n, t, b) in enumerate(steps):
    x = start_x + (col_w + gap) * i
    y = Inches(2.9)
    # column
    text(s, x, y, col_w, Inches(1.2), n, size=72, bold=True, color=ACCENT,
         font=FONT_NUM, align=PP_ALIGN.RIGHT, rtl=False)
    text(s, x, y + Inches(1.3), col_w, Inches(0.6), t, size=22, bold=True, color=INK,
         align=PP_ALIGN.RIGHT)
    text(s, x, y + Inches(2.0), col_w, Inches(2.0), b, size=14, color=MUTED,
         align=PP_ALIGN.RIGHT, line_spacing=1.55)
page_num(s, 4, TOTAL)

# ── Slide 5 — Dark mode showcase ──────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, RGBColor(0x0C, 0x0E, 0x11))
label(s, SW - Inches(2.5), Inches(0.55), "للقراءة الطويلة", color=ACCENT_2, size=11)
text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0),
     "مصمم لساعات من القراءة، لا دقائق.", size=36, bold=True,
     color=RGBColor(0xEC, 0xE6, 0xDC), align=PP_ALIGN.RIGHT)
img(s, "home_dark.png", Inches(1.2), Inches(2.1), w=Inches(11))
text(s, Inches(0.7), SH - Inches(0.55), Inches(3.5), Inches(0.4),
     "منفذ  ·  mnfz.tech", size=10, color=RGBColor(0x6D, 0x62, 0x57),
     align=PP_ALIGN.LEFT, rtl=False)
text(s, SW - Inches(1.3), SH - Inches(0.55), Inches(1), Inches(0.4),
     f"5 / {TOTAL}", size=10, color=RGBColor(0x6D, 0x62, 0x57),
     align=PP_ALIGN.RIGHT, rtl=False)

# ── Slide 6 — Features (2x2, no boxes, just text) ─────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
label(s, SW - Inches(2.5), Inches(0.7), "المزايا", color=ACCENT, size=11)
line(s, SW - Inches(2.8), Inches(1.1), Inches(0.3), Inches(0.04), color=ACCENT)
text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.8),
     "أربع مزايا.", size=44, bold=True, color=INK, align=PP_ALIGN.RIGHT)

features = [
    ("أصوات عصبية",          "Azure Neural و Kokoro — عربية طبيعية مفهومة."),
    ("تمييز كلمة بكلمة",     "كل كلمة تُضيء أثناء النطق، للقراءة الموجَّهة."),
    ("يعمل دون إنترنت",      "نسخة سطح المكتب تشغّل كل شيء محلياً، بخصوصية كاملة."),
    ("وصولية أصلية",          "وضع داكن، تكبير، اختصارات، دعم قارئ الشاشة."),
]
positions = [
    (Inches(6.95), Inches(2.9)),
    (Inches(0.7), Inches(2.9)),
    (Inches(6.95), Inches(5.0)),
    (Inches(0.7), Inches(5.0)),
]
for (t, b), (x, y) in zip(features, positions):
    line(s, x, y, Inches(0.25), Inches(0.04), color=ACCENT)
    text(s, x, y + Inches(0.15), Inches(5.7), Inches(0.6), t,
         size=22, bold=True, color=INK, align=PP_ALIGN.RIGHT)
    text(s, x, y + Inches(0.85), Inches(5.7), Inches(1.0), b,
         size=14, color=MUTED, align=PP_ALIGN.RIGHT, line_spacing=1.55)
page_num(s, 6, TOTAL)

# ── Slide 7 — Audience ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
label(s, SW - Inches(2.5), Inches(0.7), "لمن", color=ACCENT, size=11)
line(s, SW - Inches(2.8), Inches(1.1), Inches(0.3), Inches(0.04), color=ACCENT)
text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.8),
     "أربع فئات نخدمها.", size=44, bold=True, color=INK, align=PP_ALIGN.RIGHT)

groups = [
    ("مكفوفون وضعاف بصر",  "بديل عربي مجاني لقراءة الكتب والوثائق بصوت طبيعي."),
    ("ذوو عسر القراءة",     "تمييز الكلمة المنطوقة يبني جسراً بين الصوت والنص."),
    ("كبار السن",            "خط كبير، واجهة بسيطة، زر تشغيل واحد."),
    ("طلاب وباحثون",        "استمع للأبحاث والكتب أثناء التنقل أو الإجهاد البصري."),
]
positions = [
    (Inches(6.95), Inches(2.9)),
    (Inches(0.7), Inches(2.9)),
    (Inches(6.95), Inches(5.0)),
    (Inches(0.7), Inches(5.0)),
]
for (t, b), (x, y) in zip(groups, positions):
    text(s, x, y, Inches(5.7), Inches(0.6), t,
         size=26, bold=True, color=INK, align=PP_ALIGN.RIGHT)
    text(s, x, y + Inches(0.85), Inches(5.7), Inches(1.0), b,
         size=14, color=MUTED, align=PP_ALIGN.RIGHT, line_spacing=1.55)
page_num(s, 7, TOTAL)

# ── Slide 8 — Mobile + desktop showcase ───────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
label(s, SW - Inches(2.5), Inches(0.55), "حيث ما كنت", color=ACCENT, size=11)
text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0),
     "سطح المكتب، الجوال، والمتصفح.", size=36, bold=True, color=INK,
     align=PP_ALIGN.RIGHT)
# mobile on right, hero on left
img(s, "mobile_light.png", Inches(0.7), Inches(2.0), h=Inches(5.0))
img(s, "hero_light.png",   Inches(3.7), Inches(2.0), h=Inches(5.0))
text(s, Inches(0.7), SH - Inches(0.55), Inches(11.5), Inches(0.4),
     "macOS  ·  Windows  ·  Linux  ·  iOS قادم  ·  أي متصفح حديث", size=11,
     color=SOFT, align=PP_ALIGN.CENTER, rtl=False)
text(s, SW - Inches(1.3), SH - Inches(0.55), Inches(1), Inches(0.4),
     f"8 / {TOTAL}", size=10, color=SOFT, align=PP_ALIGN.RIGHT, rtl=False)

# ── Slide 9 — Roadmap (simple timeline) ───────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
label(s, SW - Inches(2.5), Inches(0.7), "الخطة", color=ACCENT, size=11)
line(s, SW - Inches(2.8), Inches(1.1), Inches(0.3), Inches(0.04), color=ACCENT)
text(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.8),
     "ماذا يُمكّنه دعم مدى.", size=44, bold=True, color=INK, align=PP_ALIGN.RIGHT)

# horizontal timeline
ty = Inches(4.2)
line(s, Inches(1), ty + Inches(0.1), Inches(11.3), Inches(0.04), color=RGBColor(0xE0, 0xD3, 0xBC))

phases = [
    ("Q4 2026", "إطلاق رسمي\nوشراكات تعليمية"),
    ("Q1 2027", "EPUB · DOCX · OCR\nللنص الممسوح"),
    ("Q2 2027", "نسخة iOS\nوأصوات عربية إضافية"),
    ("Q3 2027", "برنامج تدريب\nللمعلمين والمستخدمين"),
]
col = Inches(11.3) / 4
for i, (when, what) in enumerate(phases):
    cx = Inches(1) + col * i + col / 2
    # dot
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.12), ty - Inches(0.04), Inches(0.24), Inches(0.24))
    dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT; dot.line.fill.background()
    text(s, cx - col/2, ty - Inches(1.0), col, Inches(0.4), when,
         size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, font=FONT_NUM, rtl=False)
    text(s, cx - col/2, ty + Inches(0.55), col, Inches(1.4), what,
         size=14, color=INK, align=PP_ALIGN.CENTER, line_spacing=1.4)
page_num(s, 9, TOTAL)

# ── Slide 10 — Closing ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK); bg(s, INK)
# accent strip top
line(s, 0, 0, SW, Inches(0.18), color=ACCENT)
text(s, Inches(0), Inches(2.4), SW, Inches(1.8),
     "نحو إتاحة المعرفة\nلكل قارئ عربي.", size=58, bold=True,
     color=RGBColor(0xF7, 0xEF, 0xE2), align=PP_ALIGN.CENTER, line_spacing=1.1)
text(s, Inches(0), Inches(5.0), SW, Inches(0.5),
     "Mada Innovation Award 2026", size=12, color=ACCENT_2,
     align=PP_ALIGN.CENTER, rtl=False)
text(s, Inches(0), Inches(5.7), SW, Inches(0.5),
     "للتواصل والشراكة", size=14, color=RGBColor(0xC8, 0xBC, 0xA8),
     align=PP_ALIGN.CENTER)
text(s, Inches(0), Inches(6.25), SW, Inches(0.4),
     "oomer0171@gmail.com  ·  mnfz.tech", size=18,
     color=RGBColor(0xF7, 0xEF, 0xE2), align=PP_ALIGN.CENTER, rtl=False)

prs.save(OUT)
print(f"Saved: {OUT}")
